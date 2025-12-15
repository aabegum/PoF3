
"""
05_raporlama_ve_gorsellestirme.py (PoF3 - Ultimate Reporting Engine v3.1)

FIXES:
1. NameError: Ensures the 'charts' dictionary is explicitly defined and updated.
2. Full Visual Suite: Restores all 7 key visuals, including Aggregate Risk and Historical Trends.
"""

import os
import sys
import logging
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from datetime import datetime, timedelta

# PPTX Library Check
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Warning: 'python-pptx' not installed. Skipping PPTX generation.")

# Setup Paths
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)
from config.config import OUTPUT_DIR, LOG_DIR, INTERMEDIATE_PATHS, FEATURE_OUTPUT_PATH # Added FEATURE_OUTPUT_PATH

STEP_NAME = "05_raporlama_ve_gorsellestirme"

# Folder Setup
ACTION_DIR = os.path.join(OUTPUT_DIR, "aksiyon_listeleri")
VISUAL_DIR = os.path.join(OUTPUT_DIR, "gorseller")
REPORT_DIR = OUTPUT_DIR

for d in [ACTION_DIR, VISUAL_DIR]:
    os.makedirs(d, exist_ok=True)

# Styling
plt.style.use('ggplot')
sns.set_palette("husl")

# ------------------------------------------------------------------------------
# LOGGING
# ------------------------------------------------------------------------------
def setup_logger():
    os.makedirs(LOG_DIR, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    logger = logging.getLogger(STEP_NAME)
    logger.setLevel(logging.INFO)
    logger.handlers.clear()
    
    fh = logging.FileHandler(os.path.join(LOG_DIR, f"{STEP_NAME}_{ts}.log"), encoding="utf-8")
    fh.setFormatter(logging.Formatter("%(asctime)s - %(levelname)s - %(message)s"))
    logger.addHandler(fh)
    logger.addHandler(logging.StreamHandler(sys.stdout))
    return logger

# ------------------------------------------------------------------------------
# PHASE 1: ACTION PLANNING
# ------------------------------------------------------------------------------
def generate_action_lists(df, logger):
    logger.info("="*60)
    logger.info("[PHASE 1] Generating Action Lists...")
    
    crit_chronic = df[
        (df['Risk_Class'] == 'Critical') & 
        (df['Kronik_Seviye_Max'].isin(['KRITIK', 'YUKSEK']))
    ].copy()
    
    if not crit_chronic.empty:
        path = os.path.join(ACTION_DIR, "01_acil_mudahale_listesi.csv")
        crit_chronic.sort_values('Risk_Score', ascending=False).to_csv(path, index=False, encoding='utf-8-sig')
        logger.info(f"  > [URGENT] Chronic & Critical: {len(crit_chronic)} assets")

    trafos = df[
        (df['Ekipman_Tipi'] == 'Trafo') & 
        (df['Risk_Class'].isin(['Critical', 'High']))
    ].copy()
    
    if not trafos.empty:
        path = os.path.join(ACTION_DIR, "02_yuksek_riskli_trafolar_capex.csv")
        trafos.sort_values('Risk_Score', ascending=False).to_csv(path, index=False, encoding='utf-8-sig')
        logger.info(f"  > [CAPEX] High Risk Transformers: {len(trafos)} assets")

    inspection = df[
        (df.get('PoF_Ensemble_12Ay', 0) > 0.10) & 
        (df['Risk_Class'].isin(['Low', 'Medium']))
    ].copy()
    
    if not inspection.empty:
        path = os.path.join(ACTION_DIR, "03_bakim_rotasi_kontrol.csv")
        inspection.sort_values('PoF_Ensemble_12Ay', ascending=False).to_csv(path, index=False, encoding='utf-8-sig')
        logger.info(f"  > [OPEX] High Prob/Low Impact: {len(inspection)} assets")

    if 'Musteri_Sayisi' in df.columns:
        audit = df[
            (df['Risk_Class'] == 'Critical') & 
            ((df['Musteri_Sayisi'] <= 1) | (df['Musteri_Sayisi'].isna()))
        ].copy()
        if not audit.empty:
            path = os.path.join(ACTION_DIR, "04_veri_kalitesi_kontrol.csv")
            audit.to_csv(path, index=False, encoding='utf-8-sig')
            logger.info(f"  > [DATA] Critical Risk/Missing Data: {len(audit)} assets")

    return crit_chronic

# ------------------------------------------------------------------------------
# PHASE 2: VISUALIZATION CORE PLOTS
# ------------------------------------------------------------------------------

def plot_single_chart(df, col_x, col_y, plot_type, title, filename, logger, **kwargs):
    # EXTRACT figure size args so they don't get passed to seaborn
    width = kwargs.pop('width', 10)
    height = kwargs.pop('height', 6)
    
    plt.figure(figsize=(width, height))
    
    try:
        if plot_type == 'scatter':
            sns.scatterplot(data=df, x=col_x, y=col_y, **kwargs)
        elif plot_type == 'hist':
            sns.histplot(df[col_x], kde=True, **kwargs)
        elif plot_type == 'bar':
            sns.barplot(x=col_x, y=col_y, data=df, **kwargs)

        plt.title(title, fontsize=14)
        plt.tight_layout()
        path = os.path.join(VISUAL_DIR, filename)
        plt.savefig(path, dpi=300, bbox_inches='tight')
        plt.close()
        logger.info(f"  > Saved: {filename}")
        return path
    except Exception as e:
        logger.error(f"  [ERROR] Failed to plot {filename}: {str(e)}")
        plt.close()
        return None

def generate_visuals(df, logger):
    logger.info("="*60)
    logger.info("[PHASE 2] Generating Visual Dashboards...")
    charts = {}

# 1. RISK MATRIX
    if 'CoF_Total_Score' in df.columns and 'PoF_Ensemble_12Ay' in df.columns:
        # FIX: Robust Palette that covers English, Turkish, and Legacy labels
        robust_palette = {
            # Standard English (Current Data)
            'Critical': 'red',
            'High': 'orange', 
            'Medium': 'gold',
            'Low': 'green',
            
            # Turkish (If used later)
            'KRİTİK': 'red', 
            'YÜKSEK': 'orange',
            'ORTA': 'gold',
            'DÜŞÜK': 'green',
            
            # Legacy/Health Scores (Old Code)
            'Excellent': 'green',
            'Good': 'yellowgreen',
            'Moderate': 'orange',
            'Poor': 'red',
            
            # Fallback
            'Unknown': 'gray'
        }

        path = plot_single_chart(df, 'CoF_Total_Score', 'PoF_Ensemble_12Ay', 'scatter', 
                                 'Varlık Risk Matrisi (Risk Matrix)', "01_risk_matrisi.png", logger,
                                 hue='Risk_Class', 
                                 palette=robust_palette,  # <--- UPDATED PALETTE
                                 s=60, alpha=0.6)
        charts['risk_matrix'] = path

    # 2. HEALTH SCORE DISTRIBUTION
    if 'Health_Score' in df.columns:
        path = plot_single_chart(df, 'Health_Score', None, 'hist', 
                                 'Varlık Sağlık Skoru Dağılımı', "02_saglik_skoru_dagilimi.png", logger,
                                 bins=30, color='teal', edgecolor='black')
        charts['health_dist'] = path

    # 3. CHRONIC BREAKDOWN
    if 'Kronik_Seviye_Max' in df.columns:
        counts = df['Kronik_Seviye_Max'].value_counts()
        order = ['KRITIK', 'YUKSEK', 'ORTA', 'IZLEME', 'NORMAL']
        counts = counts.reindex([x for x in order if x in counts.index])
        
        plt.figure(figsize=(8, 6))
        colors = ['red', 'orange', 'gold', 'lightblue', 'green']
        counts.plot(kind='bar', color=colors, edgecolor='black')
        plt.title('Kronik Arıza Seviyeleri', fontsize=14)
        plt.ylabel('Ekipman Sayısı')
        plt.xticks(rotation=0)
        path = os.path.join(VISUAL_DIR, "03_kronik_dagilimi.png")
        plt.savefig(path, dpi=300, bbox_inches='tight')
        plt.close()
        charts['chronic_dist'] = path

    # 4. GEOSPATIAL MAP
    if 'Latitude' in df.columns and 'Longitude' in df.columns:
        gdf = df[(df['Latitude'] != 0) & (df['Longitude'] != 0)].copy()
        if not gdf.empty:
                # Define the palette robustly
                palette_map = {'Critical': 'red', 'Moderate': 'orange', 'Good': 'yellow', 'Excellent': 'green'}
                
                # Add fallback for missing keys if needed
                for label in gdf['Health_Class'].unique():
                    if label not in palette_map:
                        palette_map[label] = 'gray'

                path = plot_single_chart(gdf, 'Longitude', 'Latitude', 'scatter', 
                                        'Coğrafi Risk Haritası', "04_cografi_risk_haritasi.png", logger,
                                        hue='Health_Class', height=10, width=10,
                                        palette=palette_map, s=30, alpha=0.8) # Use palette_map here
                charts['geo_map'] = path

    # 5. FEATURE IMPORTANCE (Correlation Proxy)
    corr_path = os.path.join(OUTPUT_DIR, "feature_correlations.csv")
    if os.path.exists(corr_path):
        try:
            corr_df = pd.read_csv(corr_path, index_col=0)
            if 'event' in corr_df.columns:
                top_features = corr_df['event'].abs().sort_values(ascending=False).head(10).drop('event', errors='ignore')
                plt.figure(figsize=(10, 6))
                top_features.plot(kind='barh', color='purple', edgecolor='black')
                plt.title('En Önemli Risk Faktörleri (Korelasyon)', fontsize=14)
                plt.xlabel('Korelasyon Gücü')
                plt.gca().invert_yaxis()
                path = os.path.join(VISUAL_DIR, "05_ozellik_onemi.png")
                plt.savefig(path, dpi=300, bbox_inches='tight')
                plt.close()
                charts['feature_imp'] = path
        except: pass

    # 6. HISTORICAL TREND
    events_path = INTERMEDIATE_PATHS["fault_events_clean"]
    if os.path.exists(events_path):
        try:
            ev = pd.read_csv(events_path, parse_dates=['Ariza_Baslangic_Zamani'])
            ev['Year'] = ev['Ariza_Baslangic_Zamani'].dt.year
            trend = ev.groupby('Year').size()
            plt.figure(figsize=(10, 6))
            trend.plot(kind='line', marker='o', linewidth=2, color='blue')
            plt.title('Yıllık Arıza Trendi', fontsize=14)
            plt.ylabel('Toplam Arıza Sayısı')
            plt.grid(True)
            path = os.path.join(VISUAL_DIR, "06_ariza_trendi.png")
            plt.savefig(path, dpi=300, bbox_inches='tight')
            plt.close()
            charts['fault_trend'] = path
        except: pass

    # 7. AGE DISTRIBUTION
    if 'Ekipman_Yasi_Gun' in df.columns:
        ages = df['Ekipman_Yasi_Gun'] / 365.25
        path = plot_single_chart(pd.DataFrame({'Ages': ages}), 'Ages', None, 'hist', 
                                 'Varlık Yaş Dağılımı (Yıl)', "07_yas_dagilimi.png", logger,
                                 bins=20, color='brown', edgecolor='black', width=10, height=6)
        charts['age_dist'] = path

    return charts
def validate_base_rates(df, logger):
    """
    REALITY CHECK: Compares Model Predicted Failure Rates vs. Industry Standards.
    """
    logger.info("="*60)
    logger.info("[VALIDATION] Checking Model Calibration against Industry Base Rates...")
    
    # Industry Standards (The "Gold Standard" you provided)
    INDUSTRY_RANGES = {
        'Trafo':      (0.005, 0.05),  # 0.5% - 5%
        'Kesici':     (0.01, 0.08),   # 1% - 8%
        'Ayırıcı':    (0.02, 0.12),   # 2% - 12%
        'Sigorta':    (0.10, 0.40),   # 10% - 40%
        'Hat':        (0.005, 0.15),  # 0.5% - 15% (Combined OH/UG)
        'Direk':      (0.001, 0.03)   # 0.1% - 3%
    }
    
    # Normalize names if needed (Match your 'Ekipman_Tipi' values)
    # Example: 'Trafo Arızaları' -> 'Trafo'
    
    if 'PoF_Ensemble_12Ay' not in df.columns:
        logger.warning("  [SKIP] PoF column missing. Cannot validate rates.")
        return

    # Calculate Predicted Base Rate per Type
    # We count "Failure" as PoF > 0.50 (or your decision threshold)
    # OR better: Mean PoF represents the expected failure rate of the fleet
    
    stats = df.groupby('Ekipman_Tipi')['PoF_Ensemble_12Ay'].mean().reset_index()
    stats.columns = ['Type', 'Predicted_Rate']
    
    for _, row in stats.iterrows():
        etype = row['Type']
        pred = row['Predicted_Rate']
        
        # Simple fuzzy matching for keys
        matched_key = next((k for k in INDUSTRY_RANGES if k in etype), None)
        
        if matched_key:
            low, high = INDUSTRY_RANGES[matched_key]
            status = "✅ OK"
            if pred < low: status = "📉 LOW (Under-predicting)"
            if pred > high: status = "🚨 HIGH (Over-predicting)"
            
            logger.info(f"  > {etype.ljust(20)}: {pred:.1%} (Target: {low:.0%} - {high:.0%}) -> {status}")
        else:
            logger.info(f"  > {etype.ljust(20)}: {pred:.1%} (No benchmark)")

    logger.info("  [NOTE] If rates are High, consider Probability Calibration (Isotonic Regression).")
def plot_aggregate_risk_by_type(df, logger):
    """
    Plots mean PoF and Chronic Scores by Equipment Type (STOLEN FEATURE).
    """
    if 'PoF_Ensemble_12Ay' not in df.columns or 'Kronik_Kritik' not in df.columns:
        return None # Return None if data is insufficient

    agg_df = df.groupby('Ekipman_Tipi').agg(
        Mean_PoF_1Y=('PoF_Ensemble_12Ay', 'mean'),
        Mean_Chronic_Flag=('Kronik_Kritik', 'mean'),
        Count=('cbs_id', 'count')
    ).reset_index()
    
    agg_df = agg_df[agg_df['Count'] >= 100].sort_values('Mean_PoF_1Y', ascending=False).head(10)

    if agg_df.empty: return None

    fig, ax1 = plt.subplots(figsize=(12, 6))
    sns.barplot(x='Ekipman_Tipi', y='Mean_PoF_1Y', data=agg_df, ax=ax1, color='darkred', alpha=0.7)
    ax1.set_ylabel('Ortalama PoF (1 Yıl)', color='darkred', fontsize=12)
    ax1.tick_params(axis='y', labelcolor='darkred')
    ax1.set_xlabel('Ekipman Tipi', fontsize=12)
    ax1.tick_params(axis='x', rotation=45)
    
    ax2 = ax1.twinx()
    sns.lineplot(x='Ekipman_Tipi', y='Mean_Chronic_Flag', data=agg_df, ax=ax2, color='darkgreen', marker='o', linewidth=3)
    ax2.set_ylabel('Ortalama Kronik Skor (0-1)', color='darkgreen', fontsize=12)
    ax2.tick_params(axis='y', labelcolor='darkgreen')
    ax2.set_ylim(0, agg_df['Mean_Chronic_Flag'].max() * 1.2)

    plt.title('Ekipman Tipine Göre Risk Yoğunluğu (Top 10)', fontsize=14)
    
    path = os.path.join(VISUAL_DIR, "08_aggregate_risk_by_type.png")
    plt.savefig(path, dpi=300, bbox_inches='tight')
    plt.close()
    logger.info(f"  > Saved: 08_aggregate_risk_by_type.png")
    return path # Return the path

# ------------------------------------------------------------------------------
# PHASE 3: EXCEL REPORTING
# ------------------------------------------------------------------------------
# Updated signature to accept 'case_studies'
def create_excel_report(df, crit_chronic, case_studies, logger): 
    logger.info("="*60)
    logger.info("[PHASE 3] Creating Excel Report...")
    
    timestamp = datetime.now().strftime("%Y-%m-%d")
    out_path = os.path.join(REPORT_DIR, f"PoF3_Analiz_Raporu_Final.xlsx")
    
    with pd.ExcelWriter(out_path, engine='openpyxl') as writer:
        # 1. Summary (YOUR ORIGINAL LOGIC KEPT INT ACT)
        total = len(df)
        crit_count = (df['Risk_Class'] == 'Critical').sum() if 'Risk_Class' in df.columns else 0
        avg_health = df['Health_Score'].mean() if 'Health_Score' in df.columns else 0
        
        summary = pd.DataFrame({
            'KPI': ['Toplam Varlık', 'Kritik Riskli', 'Kronik ve Kritik', 'Ortalama Sağlık', 'Rapor Tarihi'],
            'Değer': [total, crit_count, len(crit_chronic), f"{avg_health:.1f}", timestamp]
        })
        summary.to_excel(writer, sheet_name='Yonetici_Ozeti', index=False)
        
        # 2. Action List (YOUR ORIGINAL LOGIC)
        if not crit_chronic.empty:
            crit_chronic.to_excel(writer, sheet_name='Acil_Mudahale', index=False)

        # --- NEW: Case Studies Tab ---
        if not case_studies.empty:
            case_studies.to_excel(writer, sheet_name='Vaka_Analizi_CaseStudy', index=False)
        # -----------------------------
            
        # 3. Top 1000 Master (YOUR ORIGINAL LOGIC)
        if 'Risk_Score' in df.columns:
            df.sort_values('Risk_Score', ascending=False).head(1000).to_excel(writer, sheet_name='Risk_Master_Top1000', index=False)
        else:
            df.head(1000).to_excel(writer, sheet_name='Risk_Master_Top1000', index=False)
            
    logger.info(f"  > Saved: {os.path.basename(out_path)}")
    
def generate_case_studies(df_risk, logger):
    """
    STOLEN FEATURE #3: Case Study Generator.
    Identifies assets that failed recently (last 6 months) and checks if
    the model correctly flagged them as High/Critical Risk.
    """
    logger.info("[PHASE 1.5] Generating Failure Case Studies...")
    
    # 1. Load Raw Events
    events_path = INTERMEDIATE_PATHS["fault_events_clean"]
    if not os.path.exists(events_path):
        logger.warning(f"  [WARN] Event file not found at {events_path}. Skipping Case Studies.")
        return pd.DataFrame()
        
    events = pd.read_csv(events_path, parse_dates=['Ariza_Baslangic_Zamani'])
    events['cbs_id'] = events['cbs_id'].astype(str).str.lower().str.strip()
    
    # 2. Filter for Recent Failures (Last 180 Days)
    if events.empty:
        return pd.DataFrame()

    analysis_date = events['Ariza_Baslangic_Zamani'].max()
    cutoff_date = analysis_date - timedelta(days=180)
    recent_faults = events[events['Ariza_Baslangic_Zamani'] >= cutoff_date].copy()
    
    if recent_faults.empty:
        logger.info("  [INFO] No recent faults found within 180 days. Skipping Case Studies.")
        return pd.DataFrame()

    # 3. Merge with Risk Predictions
    # FIX: We explicitly handle suffixes to avoid 'Ekipman_Tipi' collision
    # We bring in 'Ekipman_Tipi' from risk data as 'Ekipman_Tipi_Risk' just in case
    cols_to_merge = ['cbs_id', 'Risk_Class', 'Risk_Score', 'PoF_Ensemble_12Ay', 'Ilce']
    
    # If df_risk has Ekipman_Tipi, rename it temporarily or rely on suffixes
    if 'Ekipman_Tipi' in df_risk.columns:
        cols_to_merge.append('Ekipman_Tipi')

    case_df = recent_faults.merge(
        df_risk[cols_to_merge], 
        on='cbs_id', 
        how='left',
        suffixes=('', '_RiskContext') # Left keeps original name, Right gets suffix
    )
    
    # 4. Column Cleanup (The Fix for KeyError)
    # If collision happened, 'Ekipman_Tipi' (from events) exists. 
    # If 'Ekipman_Tipi' was only in Risk DF, it might be named 'Ekipman_Tipi_RiskContext'.
    
    if 'Ekipman_Tipi' not in case_df.columns and 'Ekipman_Tipi_RiskContext' in case_df.columns:
        case_df.rename(columns={'Ekipman_Tipi_RiskContext': 'Ekipman_Tipi'}, inplace=True)
    elif 'Ekipman_Tipi' in case_df.columns and 'Ekipman_Tipi_RiskContext' in case_df.columns:
        # Fill missing event types with risk master types
        case_df['Ekipman_Tipi'] = case_df['Ekipman_Tipi'].fillna(case_df['Ekipman_Tipi_RiskContext'])

    # 5. Classify Outcomes
    def judge_prediction(row):
        if pd.isna(row['Risk_Class']): return "Unknown Asset"
        if row['Risk_Class'] in ['Critical', 'High', 'KRİTİK', 'YÜKSEK']: # Handle Turkish labels too
            return "SUCCESS (Predicted)"
        elif row['Risk_Class'] in ['Medium', 'ORTA']:
            return "WATCHLIST (Partial)"
        else:
            return "MISS (Low Risk)"

    case_df['Model_Verdict'] = case_df.apply(judge_prediction, axis=1)
    
    # 6. Select Top Examples
    successes = case_df[case_df['Model_Verdict'] == "SUCCESS (Predicted)"].sort_values('Risk_Score', ascending=False).head(10)
    misses = case_df[case_df['Model_Verdict'] == "MISS (Low Risk)"].sort_values('Risk_Score', ascending=True).head(5)
    
    final_cases = pd.concat([successes, misses])
    
    # Final Safety Check for Columns
    report_cols = ['cbs_id', 'Ariza_Baslangic_Zamani', 'Ekipman_Tipi', 'Ilce', 'Risk_Class', 'Risk_Score', 'Model_Verdict']
    available_cols = [c for c in report_cols if c in final_cases.columns]
    
    final_cases = final_cases[available_cols].sort_values('Risk_Score', ascending=False)
    
    logger.info(f"  > Generated {len(final_cases)} case studies.")
    return final_cases
    # 4. Classify Outcomes
    def judge_prediction(row):
        if pd.isna(row['Risk_Class']): return "Unknown Asset"
        if row['Risk_Class'] in ['Critical', 'High']:
            return "SUCCESS (Predicted)"
        elif row['Risk_Class'] == 'Medium':
            return "WATCHLIST (Partial)"
        else:
            return "MISS (Low Risk)"

    case_df['Model_Verdict'] = case_df.apply(judge_prediction, axis=1)
    
    # 5. Select Top Examples
    # Get top 10 Successes (Highest Risk Score) and top 5 Misses (Lowest Risk Score)
    successes = case_df[case_df['Model_Verdict'] == "SUCCESS (Predicted)"].sort_values('Risk_Score', ascending=False).head(10)
    misses = case_df[case_df['Model_Verdict'] == "MISS (Low Risk)"].sort_values('Risk_Score', ascending=True).head(5)
    
    final_cases = pd.concat([successes, misses])
    
    # Clean up for Report
    report_cols = ['cbs_id', 'Ariza_Baslangic_Zamani', 'Ekipman_Tipi', 'Ilce', 'Risk_Class', 'Risk_Score', 'Model_Verdict']
    final_cases = final_cases[report_cols].sort_values('Risk_Score', ascending=False)
    
    logger.info(f"  > Generated {len(final_cases)} case studies (Successes vs Misses).")
    return final_cases
# ------------------------------------------------------------------------------
# PHASE 4: POWERPOINT PRESENTATION (RESTORED)
# ------------------------------------------------------------------------------
def create_pptx_presentation(df, charts, logger):
    if not HAS_PPTX:
        logger.warning("[PHASE 4] Skipping PPTX (Library missing).")
        return

    logger.info("="*60)
    logger.info("[PHASE 4] Creating PowerPoint Dashboard...")
    
    prs = Presentation()
    timestamp = datetime.now().strftime("%d %B %Y")
    
    # 1. Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "PoF3 Risk ve Sağlık Analizi"
    slide.placeholders[1].text = f"Yönetici Özeti Raporu\n{timestamp}"
    
    # 2. Executive Summary Slide
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Genel Durum Özeti"
    
    total_assets = len(df)
    critical_assets = (df['Risk_Class'] == 'Critical').sum() if 'Risk_Class' in df.columns else 0
    chronic_assets = (df['Kronik_Seviye_Max'] != 'NORMAL').sum() if 'Kronik_Seviye_Max' in df.columns else 0
    avg_health = df['Health_Score'].mean() if 'Health_Score' in df.columns else 0
    
    content = slide.placeholders[1]
    content.text = (
        f"Toplam Varlık Sayısı: {total_assets:,}\n"
        f"Kritik Riskli Varlıklar: {critical_assets:,} ({(critical_assets/total_assets):.1%})\n"
        f"Kronik Sorunlu Varlıklar: {chronic_assets:,}\n"
        f"Filo Ortalama Sağlık Skoru: {avg_health:.1f} / 100\n\n"
        "Öneri: 'Kritik' ve 'Kronik' kesişimindeki varlıklara öncelik verilmelidir."
    )

    # 3. Add Charts
    # Map chart names to Slide Titles
    chart_slides = {
        'risk_matrix': "Risk Matrisi (Etki vs Olasılık)",
        'health_dist': "Filo Sağlık Dağılımı",
        'chronic_dist': "Kronik Arıza Analizi",
        'aggregate_risk': "Ekipman Tipine Göre Risk Yoğunluğu", # NEW AGGREGATE PLOT
        'feature_imp': "Risk Faktörleri (Önem Derecesi)",
        'fault_trend': "Tarihsel Arıza Trendi",
        'age_dist': "Varlık Yaş Profili",
        'geo_map': "Coğrafi Risk Haritası"
    }
    
    for key, slide_title in chart_slides.items():
        if key in charts and os.path.exists(charts[key]):
            slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
            slide.shapes.title.text = slide_title
            
            # Add Image (Centered)
            left = Inches(1)
            top = Inches(1.5)
            height = Inches(5.5)
            slide.shapes.add_picture(charts[key], left, top, height=height)
            
    out_path = os.path.join(OUTPUT_DIR, "PoF3_Yonetici_Sunumu_Final.pptx")
    prs.save(out_path)
    logger.info(f"  > Saved: {os.path.basename(out_path)}")

# ------------------------------------------------------------------------------
# MAIN
# ------------------------------------------------------------------------------
def main():
    logger = setup_logger()
    
    # 1. Load Risk Data
    risk_path = os.path.join(OUTPUT_DIR, "risk_equipment_master.csv")
    if not os.path.exists(risk_path):
        logger.error(f"[FATAL] Risk Master not found at {risk_path}. Run Step 04 first.")
        return
        
    df = pd.read_csv(risk_path)
    
    # 2. Load & Merge Context (Geo + Ilce/Sehir)
    master_path = os.path.join(os.path.dirname(OUTPUT_DIR), "ara_ciktilar", "equipment_master.csv")
    if os.path.exists(master_path):
        meta = pd.read_csv(master_path)
        meta['cbs_id'] = meta['cbs_id'].astype(str).str.lower().str.strip()
        
        # Identify columns to bring over (Geo + Admin Locations)
        desired_cols = ['Latitude', 'Longitude', 'Musteri_Sayisi', 'Ilce', 'Sehir', 'Mahalle', 'Ekipman_Tipi']
        cols_to_add = [c for c in desired_cols if c in meta.columns and c not in df.columns]
        
        if cols_to_add:
            logger.info(f"[MERGE] Adding context columns: {cols_to_add}")
            df['cbs_id'] = df['cbs_id'].astype(str).str.lower().str.strip()
            df = df.merge(meta[['cbs_id'] + cols_to_add], on='cbs_id', how='left')
    
    # 3. Safety Fill (Prevents 'KeyError')
    # If Ilce still doesn't exist (e.g., merge failed), create it.
    for col in ['Ilce', 'Sehir', 'Ekipman_Tipi']:
        if col not in df.columns:
            df[col] = 'Unknown'
        else:
            df[col] = df[col].fillna('Unknown')

    logger.info(f"[LOAD] Reporting on {len(df):,} assets.")
    validate_base_rates(df, logger)
    # 4. Execute Pipeline Stages
    crit_chronic = generate_action_lists(df, logger)
    
    # Case Studies (Safe to run now that Ilce exists)
    case_studies = generate_case_studies(df, logger)
    
    charts = generate_visuals(df, logger)
    
    # Aggregate Plot
    agg_path = plot_aggregate_risk_by_type(df, logger)
    if agg_path:
        charts['aggregate_risk'] = agg_path 
    
    # 5. Generate Final Outputs
    create_excel_report(df, crit_chronic, case_studies, logger)
    create_pptx_presentation(df, charts, logger)
    
    logger.info("")
    logger.info("[SUCCESS] Reporting & Visualization Complete.")

if __name__ == "__main__":
    main()


if __name__ == "__main__":
    main()
