"""
07_generate_deliverables_v3.py - PoF3 Professional Output Package Generator (v3)

SOLVES THE PROBLEM:
- 20+ scattered CSVs → 1 Excel workbook + 1 PowerPoint
- Multiple models → Unified PoF view (Cox, Weibull, RSF, ML, ENSEMBLE)
- No executive summary → Management-ready dashboard

KEY INPUTS (from PoF3 v3 pipeline):
- data/ara_ciktilar/equipment_master.csv
- data/ara_ciktilar/ozellikler_pof3.csv
- data/sonuclar/risk_equipment_master.csv (if available)
- data/sonuclar/chronic_equipment_summary.csv
- data/sonuclar/temporal_stability_report.csv (optional)
- data/ara_ciktilar/data_range_metadata.csv
- data/sonuclar/cox_sagkalim_*ay_ariza_olasiligi.csv
- data/sonuclar/weibull_sagkalim_*ay_ariza_olasiligi.csv
- data/sonuclar/rsf_sagkalim_*ay_ariza_olasiligi.csv
- data/sonuclar/ensemble_sagkalim_*ay_ariza_olasiligi.csv
- data/sonuclar/leakage_free_ml_pof.csv
- data/sonuclar/shap_feature_importance.csv (optional)
- data/sonuclar/model_comparison_report.csv (optional)

OUTPUTS:
1. PoF_Analysis_YYYY-MM-DD_HHMM.xlsx
   - Executive Summary
   - High Risk Top 100 (Ensemble 12 Ay)
   - Risk Equipment Master (if available)
   - PoF Predictions (all models, all horizons)
   - Chronic Equipment (Step 04 summary)
   - Model Comparison (optional)
   - Temporal Analysis (optional)
   - Feature Importance (optional)
   - Metadata

2. PoF_Dashboard_YYYY-MM-DD_HHMM.pptx
   - Yönetici Özeti
   - Equipment Distribution
   - Chronic Level Distribution
   - PoF Distribution by Horizon (Ensemble / Cox)
   - Feature Importance (SHAP)
   - Temporal Trends

3. High-quality PNG charts in gorseller/ folder

READY FOR:
- EDAŞ management presentations
- Technical committees
- Field operation planning
"""

import os
import sys
from pathlib import Path
from datetime import datetime
import warnings

warnings.filterwarnings("ignore")

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns

# Excel / PowerPoint libraries
try:
    from openpyxl import load_workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    HAS_OPENPYXL = True
except Exception:
    HAS_OPENPYXL = False
    print("Warning: openpyxl not installed. Install: pip install openpyxl")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except Exception:
    HAS_PPTX = False
    print("Warning: python-pptx not installed. Install: pip install python-pptx")

# ---------------------------------------------------------------------------
# PATHS / GLOBALS
# ---------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parents[1]
DATA_DIR = ROOT / "data"
RESULTS_DIR = DATA_DIR / "sonuclar"
INTER_DIR = DATA_DIR / "ara_ciktilar"
VIS_DIR = ROOT / "gorseller"
VIS_DIR.mkdir(parents=True, exist_ok=True)

timestamp = datetime.now().strftime("%Y-%m-%d_%H%M")
EXCEL_OUTPUT = RESULTS_DIR / f"PoF_Analysis_{timestamp}.xlsx"
PPTX_OUTPUT = RESULTS_DIR / f"PoF_Dashboard_{timestamp}.pptx"

# Plot styling
sns.set_style("whitegrid")
plt.rcParams["figure.dpi"] = 150
plt.rcParams["font.family"] = "DejaVu Sans"
plt.rcParams["axes.unicode_minus"] = False

# EDAŞ-friendly color palette
COLORS = {
    "primary": "#1f77b4",    # Blue
    "danger": "#d62728",     # Red
    "warning": "#ff7f0e",    # Orange
    "success": "#2ca02c",    # Green
    "info": "#17becf",       # Cyan
    "kritik": "#d62728",     # Red
    "yuksek": "#ff7f0e",     # Orange
    "orta": "#ffbb00",       # Yellow/orange
    "izleme": "#17becf",     # Cyan
    "normal": "#2ca02c",     # Green,
}


# ======================================================================
# UTILS
# ======================================================================
def safe_read_csv(path: Path, name: str, encoding="utf-8-sig") -> pd.DataFrame:
    """Safe CSV reader with logging."""
    if path.exists():
        try:
            df = pd.read_csv(path, encoding=encoding)
            print(f"✓ Loaded {name}: {len(df):,} rows from {path.name}")
            return df
        except Exception as e:
            print(f"✗ Failed to load {name} from {path}: {e}")
            return pd.DataFrame()
    else:
        print(f"✗ Missing {name}: {path}")
        return pd.DataFrame()


def load_metadata() -> dict:
    """Load DATA_START_DATE / DATA_END_DATE / ANALYSIS_DATE from metadata."""
    meta = {
        "DATA_START_DATE": None,
        "DATA_END_DATE": None,
        "ANALYSIS_DATE": None,
    }

    metadata_path = INTER_DIR / "data_range_metadata.csv"
    if not metadata_path.exists():
        print(f"✗ Metadata file not found: {metadata_path} (using defaults)")
        return meta

    try:
        df = pd.read_csv(metadata_path, encoding="utf-8-sig")
        for param in ["DATA_START_DATE", "DATA_END_DATE", "ANALYSIS_DATE"]:
            if param in df["Parameter"].values:
                value = df.loc[df["Parameter"] == param, "Value"].iloc[0]
                meta[param] = str(value)
        print(f"✓ Loaded metadata from {metadata_path.name}: {meta}")
    except Exception as e:
        print(f"✗ Failed to parse metadata: {e}")

    return meta


# ======================================================================
# DATA LOADING
# ======================================================================
def load_all_data():
    """Load all necessary datasets for deliverables."""
    print("\n" + "=" * 80)
    print("🚀 LOADING POF3 DATA (v3)")
    print("=" * 80)

    data = {}

    # Metadata
    data["meta"] = load_metadata()

    # Core datasets
    data["equipment"] = safe_read_csv(
        INTER_DIR / "equipment_master.csv", "equipment_master"
    )
    data["features"] = safe_read_csv(
        INTER_DIR / "ozellikler_pof3.csv", "ozellikler_pof3"
    )
    data["chronic"] = safe_read_csv(
        RESULTS_DIR / "chronic_equipment_summary.csv", "chronic_equipment_summary"
    )
    data["temporal"] = safe_read_csv(
        RESULTS_DIR / "temporal_stability_report.csv", "temporal_stability_report"
    )

    # Risk master (preferred for most downstream views)
    risk_paths = [
        RESULTS_DIR / "risk_equipment_master.csv",
        INTER_DIR / "risk_equipment_master.csv",
    ]
    risk_df = pd.DataFrame()
    for p in risk_paths:
        if p.exists():
            risk_df = safe_read_csv(p, "risk_equipment_master")
            break
    data["risk_master"] = risk_df

    # PoF predictions
    data["pof"] = {}
    pof_files = {
        # Cox calibrated (multi-horizon)
        "cox_3mo": RESULTS_DIR / "cox_sagkalim_3ay_ariza_olasiligi.csv",
        "cox_6mo": RESULTS_DIR / "cox_sagkalim_6ay_ariza_olasiligi.csv",
        "cox_12mo": RESULTS_DIR / "cox_sagkalim_12ay_ariza_olasiligi.csv",
        "cox_24mo": RESULTS_DIR / "cox_sagkalim_24ay_ariza_olasiligi.csv",
        # Weibull AFT (optional)
        "weibull_3mo": RESULTS_DIR / "weibull_sagkalim_3ay_ariza_olasiligi.csv",
        "weibull_6mo": RESULTS_DIR / "weibull_sagkalim_6ay_ariza_olasiligi.csv",
        "weibull_12mo": RESULTS_DIR / "weibull_sagkalim_12ay_ariza_olasiligi.csv",
        "weibull_24mo": RESULTS_DIR / "weibull_sagkalim_24ay_ariza_olasiligi.csv",
        # RSF (multi-horizon if available)
        "rsf_3mo": RESULTS_DIR / "rsf_sagkalim_3ay_ariza_olasiligi.csv",
        "rsf_6mo": RESULTS_DIR / "rsf_sagkalim_6ay_ariza_olasiligi.csv",
        "rsf_12mo": RESULTS_DIR / "rsf_sagkalim_12ay_ariza_olasiligi.csv",
        "rsf_24mo": RESULTS_DIR / "rsf_sagkalim_24ay_ariza_olasiligi.csv",
        # Ensemble (RECOMMENDED)
        "ensemble_3mo": RESULTS_DIR / "ensemble_sagkalim_3ay_ariza_olasiligi.csv",
        "ensemble_6mo": RESULTS_DIR / "ensemble_sagkalim_6ay_ariza_olasiligi.csv",
        "ensemble_12mo": RESULTS_DIR / "ensemble_sagkalim_12ay_ariza_olasiligi.csv",
        "ensemble_24mo": RESULTS_DIR / "ensemble_sagkalim_24ay_ariza_olasiligi.csv",
        # ML leakage-free PoF
        "ml": RESULTS_DIR / "leakage_free_ml_pof.csv",
    }

    for key, path in pof_files.items():
        df = safe_read_csv(path, key)
        if not df.empty:
            data["pof"][key] = df

    # Model comparison report (optional)
    data["comparison"] = safe_read_csv(
        RESULTS_DIR / "model_comparison_report.csv", "model_comparison_report"
    )

    # Feature importance (SHAP)
    data["shap"] = safe_read_csv(
        RESULTS_DIR / "shap_feature_importance.csv", "shap_feature_importance"
    )

    print("\n✓ Data loading completed.")
    return data


# ======================================================================
# VISUALIZATION GENERATION
# ======================================================================
def create_visualizations(data):
    """Generate all visualizations used in Excel + PowerPoint."""
    print("\n" + "=" * 80)
    print("🎨 GENERATING VISUALIZATIONS")
    print("=" * 80)

    charts = {}

    # Decide base equipment DF (risk_master preferred)
    base_df = (
        data["risk_master"]
        if not data["risk_master"].empty
        else (data["features"] if not data["features"].empty else data["equipment"])
    )

    # 1) Equipment Type Distribution
    if base_df is not None and not base_df.empty and "Ekipman_Tipi" in base_df.columns:
        fig, ax = plt.subplots(figsize=(10, 6))
        counts = base_df["Ekipman_Tipi"].value_counts().head(10)
        counts.plot.barh(ax=ax, color=COLORS["primary"])
        ax.set_xlabel("Ekipman Sayısı", fontsize=12)
        ax.set_title("Ekipman Tip Dağılımı (Top 10)", fontsize=14, fontweight="bold")
        plt.tight_layout()
        chart_path = VIS_DIR / "equipment_distribution.png"
        plt.savefig(chart_path, dpi=300, bbox_inches="tight")
        plt.close()
        charts["equipment_dist"] = chart_path
        print("✓ Chart: equipment_distribution.png")

    # 2) Chronic Severity Distribution (Step 04 v3, with IZLEME)
    if not data["chronic"].empty and "Kronik_Seviye_Max" in data["chronic"].columns:
        fig, ax = plt.subplots(figsize=(10, 6))
        sev_counts = data["chronic"]["Kronik_Seviye_Max"].value_counts()
        sev_order = ["KRITIK", "YUKSEK", "ORTA", "IZLEME", "NORMAL"]
        sev_counts = sev_counts.reindex(
            [s for s in sev_order if s in sev_counts.index]
        )

        color_map = {
            "KRITIK": COLORS["kritik"],
            "YUKSEK": COLORS["yuksek"],
            "ORTA": COLORS["orta"],
            "IZLEME": COLORS["izleme"],
            "NORMAL": COLORS["normal"],
        }
        colors = [color_map.get(s, COLORS["info"]) for s in sev_counts.index]

        sev_counts.plot.bar(ax=ax, color=colors)
        ax.set_xlabel("Kronik Seviye", fontsize=12)
        ax.set_ylabel("Ekipman Sayısı", fontsize=12)
        ax.set_title("Kronik Ekipman Dağılımı (IEEE 1366 Uyumlu)", fontsize=14, fontweight="bold")
        plt.xticks(rotation=0)
        plt.tight_layout()
        chart_path = VIS_DIR / "chronic_distribution.png"
        plt.savefig(chart_path, dpi=300, bbox_inches="tight")
        plt.close()
        charts["chronic_dist"] = chart_path
        print("✓ Chart: chronic_distribution.png")

    # 3) PoF Distribution by Horizon (prefer ENSEMBLE, fallback COX)
    if data["pof"]:
        fig, axes = plt.subplots(2, 2, figsize=(14, 10))
        axes = axes.flatten()

        configs = [
            ("3", "3 Ay"),
            ("6", "6 Ay"),
            ("12", "12 Ay"),
            ("24", "24 Ay"),
        ]

        for idx, (m, label) in enumerate(configs):
            ax = axes[idx]
            key_ens = f"ensemble_{m}mo"
            key_cox = f"cox_{m}mo"

            df = None
            model_label = ""
            if key_ens in data["pof"]:
                df = data["pof"][key_ens]
                model_label = "Ensemble"
            elif key_cox in data["pof"]:
                df = data["pof"][key_cox]
                model_label = "Cox (kalibre)"

            if df is None or df.empty:
                ax.set_title(f"PoF Dağılımı - {label} (veri yok)")
                ax.axis("off")
                continue

            pof_cols = [c for c in df.columns if c != "cbs_id"]
            if not pof_cols:
                ax.set_title(f"PoF Dağılımı - {label} (geçersiz format)")
                ax.axis("off")
                continue

            pofs = df[pof_cols[0]].dropna()
            ax.hist(
                pofs,
                bins=50,
                color=COLORS["primary"],
                alpha=0.7,
                edgecolor="black",
            )
            ax.set_title(
                f"PoF Dağılımı - {label} ({model_label})",
                fontsize=12,
                fontweight="bold",
            )
            ax.set_xlabel("PoF", fontsize=10)
            ax.set_ylabel("Ekipman Sayısı", fontsize=10)
            median = pofs.median()
            ax.axvline(
                median,
                color=COLORS["danger"],
                linestyle="--",
                label=f"Median: {median:.3f}",
            )
            ax.legend()

        plt.tight_layout()
        chart_path = VIS_DIR / "pof_by_horizon.png"
        plt.savefig(chart_path, dpi=300, bbox_inches="tight")
        plt.close()
        charts["pof_horizons"] = chart_path
        print("✓ Chart: pof_by_horizon.png")

    # 4) Feature Importance (SHAP for ML)
    if not data["shap"].empty:
        fig, ax = plt.subplots(figsize=(10, 8))
        # Expect columns: feature, shap_importance
        df_shap = data["shap"].copy()
        if df_shap.shape[1] >= 2:
            df_shap = df_shap.sort_values(df_shap.columns[1], ascending=False).head(15)
            ax.barh(
                range(len(df_shap)),
                df_shap.iloc[:, 1],
                color=COLORS["success"],
            )
            ax.set_yticks(range(len(df_shap)))
            ax.set_yticklabels(df_shap.iloc[:, 0], fontsize=9)
            ax.set_xlabel(df_shap.columns[1], fontsize=12)
            ax.set_title("Top 15 Önemli Özellik (SHAP - ML PoF)", fontsize=14, fontweight="bold")
            ax.invert_yaxis()
            plt.tight_layout()
            chart_path = VIS_DIR / "feature_importance.png"
            plt.savefig(chart_path, dpi=300, bbox_inches="tight")
            plt.close()
            charts["feature_imp"] = chart_path
            print("✓ Chart: feature_importance.png")

    # 5) Temporal Fault Trends
    if not data["temporal"].empty and {
        "Year",
        "Total_Faults",
    }.issubset(data["temporal"].columns):
        fig, ax = plt.subplots(figsize=(12, 6))
        tmp = data["temporal"].copy()
        ax.plot(
            tmp["Year"],
            tmp["Total_Faults"],
            marker="o",
            linewidth=2,
            markersize=6,
            color=COLORS["primary"],
        )
        ax.fill_between(
            tmp["Year"],
            tmp["Total_Faults"],
            alpha=0.2,
            color=COLORS["primary"],
        )
        ax.set_xlabel("Yıl", fontsize=12)
        ax.set_ylabel("Toplam Arıza Sayısı", fontsize=12)
        ax.set_title("Yıllara Göre Arıza Trendi", fontsize=14, fontweight="bold")
        ax.grid(True, alpha=0.3)
        plt.tight_layout()
        chart_path = VIS_DIR / "fault_trends.png"
        plt.savefig(chart_path, dpi=300, bbox_inches="tight")
        plt.close()
        charts["fault_trends"] = chart_path
        print("✓ Chart: fault_trends.png")

    print(f"\n✓ Generated {len(charts)} visualizations.")
    return charts


# ======================================================================
# EXCEL WORKBOOK
# ======================================================================
def create_excel_workbook(data, charts):
    if not HAS_OPENPYXL:
        print("\n✗ Skipping Excel generation (openpyxl not installed).")
        return None

    print("\n" + "=" * 80)
    print("📊 GENERATING EXCEL WORKBOOK")
    print("=" * 80)

    writer = pd.ExcelWriter(EXCEL_OUTPUT, engine="openpyxl")

    meta = data.get("meta", {})
    data_start = meta.get("DATA_START_DATE") or "N/A"
    data_end = meta.get("DATA_END_DATE") or "N/A"
    analysis_date = meta.get("ANALYSIS_DATE") or datetime.now().strftime("%Y-%m-%d")

    # Base DF (risk_master preferred)
    base_df = (
        data["risk_master"]
        if not data["risk_master"].empty
        else (data["features"] if not data["features"].empty else data["equipment"])
    )

    # Sheet 1: Executive Summary
    total_eq = len(base_df) if base_df is not None and not base_df.empty else 0

    chronic_eq = 0
    if not data["chronic"].empty and "Kronik_Seviye_Max" in data["chronic"].columns:
        chronic_eq = (data["chronic"]["Kronik_Seviye_Max"] != "NORMAL").sum()

    summary_data = {
        "Metric": [
            "Toplam Ekipman",
            "Kronik Ekipman (Herhangi Seviye)",
            "KRITIK Kronik",
            "YUKSEK Kronik",
            "ORTA Kronik",
            "Veri Başlangıç Tarihi",
            "Veri Bitiş Tarihi",
            "Analiz Tarihi",
            "Önerilen Risk Modeli",
        ],
        "Value": [
            f"{total_eq:,}",
            f"{chronic_eq:,} ({(100*chronic_eq/total_eq):.1f}%)"
            if total_eq > 0
            else "0 (0.0%)",
            data["chronic"]["Kronik_Kritik"].sum()
            if not data["chronic"].empty and "Kronik_Kritik" in data["chronic"].columns
            else 0,
            data["chronic"]["Kronik_Yuksek"].sum()
            if not data["chronic"].empty and "Kronik_Yuksek" in data["chronic"].columns
            else 0,
            data["chronic"]["Kronik_Orta"].sum()
            if not data["chronic"].empty and "Kronik_Orta" in data["chronic"].columns
            else 0,
            data_start,
            data_end,
            analysis_date,
            "ENSEMBLE 12 Ay (Cox + Weibull + RSF + ML)",
        ],
    }
    summary_df = pd.DataFrame(summary_data)
    summary_df.to_excel(writer, sheet_name="Executive Summary", index=False)
    print("✓ Sheet: Executive Summary")

    # Sheet 2: High Risk Top 100 (Ensemble 12m)
    if (
        base_df is not None
        and not base_df.empty
        and "ensemble_12mo" in data["pof"]
        and not data["pof"]["ensemble_12mo"].empty
    ):
        df_ens = data["pof"]["ensemble_12mo"]
        pof_cols = [c for c in df_ens.columns if c != "cbs_id"]
        if pof_cols:
            ens_col = pof_cols[0]
            high_risk = base_df.merge(
                df_ens[["cbs_id", ens_col]], on="cbs_id", how="left"
            )
            high_risk = high_risk.sort_values(ens_col, ascending=False).head(100)

            cols_pref = [
                "cbs_id",
                "Ekipman_Tipi",
                ens_col,
                "Kronik_Seviye_Max",
                "Faults_Last_365d",
                "Ekipman_Yasi_Gun",
                "Gerilim_Seviyesi",
                "Marka",
                "Lambda_Yillik_Ariza",
            ]
            cols_final = [c for c in cols_pref if c in high_risk.columns]
            high_risk[cols_final].to_excel(
                writer, sheet_name="High Risk Top 100", index=False
            )
            print("✓ Sheet: High Risk Top 100")

    # Sheet 3: Risk Equipment Master / Features
    if base_df is not None and not base_df.empty:
        sheet_name = (
            "Risk Equipment Master"
            if not data["risk_master"].empty
            else "Equipment Master"
        )
        base_df.to_excel(writer, sheet_name=sheet_name, index=False)
        print(f"✓ Sheet: {sheet_name}")

    # Sheet 4: PoF Predictions (all models & horizons)
    if base_df is not None and not base_df.empty and data["pof"]:
        pof_combined = base_df[["cbs_id"]].drop_duplicates().copy()

        for key, df in data["pof"].items():
            if df.empty or "cbs_id" not in df.columns:
                continue

            # ML file may contain multiple probability columns
            if key == "ml":
                pof_cols = [c for c in df.columns if c != "cbs_id"]
                for col in pof_cols:
                    tmp = df[["cbs_id", col]].copy()
                    pof_combined = pof_combined.merge(tmp, on="cbs_id", how="left")
            else:
                pof_cols = [c for c in df.columns if c != "cbs_id"]
                if not pof_cols:
                    continue
                col = pof_cols[0]
                tmp = df[["cbs_id", col]].copy()
                pof_combined = pof_combined.merge(tmp, on="cbs_id", how="left")

        pof_combined.to_excel(writer, sheet_name="PoF Predictions", index=False)
        print("✓ Sheet: PoF Predictions")

    # Sheet 5: Chronic Equipment Detail
    if not data["chronic"].empty:
        # Only non-normal
        chronic_only = data["chronic"][
            data["chronic"].get("Kronik_Seviye_Max", "NORMAL") != "NORMAL"
        ]
        chronic_only.to_excel(writer, sheet_name="Chronic Equipment", index=False)
        print("✓ Sheet: Chronic Equipment")

    # Sheet 6: Model Comparison
    if not data["comparison"].empty:
        data["comparison"].to_excel(writer, sheet_name="Model Comparison", index=False)
        print("✓ Sheet: Model Comparison")

    # Sheet 7: Temporal Analysis
    if not data["temporal"].empty:
        data["temporal"].to_excel(writer, sheet_name="Temporal Analysis", index=False)
        print("✓ Sheet: Temporal Analysis")

    # Sheet 8: Feature Importance
    if not data["shap"].empty:
        data["shap"].to_excel(writer, sheet_name="Feature Importance", index=False)
        print("✓ Sheet: Feature Importance")

    # Sheet 9: Metadata
    meta_rows = []
    meta_rows.append(
        {"Parameter": "Pipeline Run Date", "Value": datetime.now().strftime("%Y-%m-%d %H:%M")}
    )
    meta_rows.append(
        {
            "Parameter": "Data Range",
            "Value": f"{data_start} to {data_end}",
        }
    )
    meta_rows.append(
        {
            "Parameter": "Total Equipment",
            "Value": total_eq,
        }
    )
    meta_rows.append(
        {
            "Parameter": "Models Used",
            "Value": "Cox PH (kalibre), Weibull AFT, RSF, XGBoost ML, ENSEMBLE (önerilen)",
        }
    )
    metadata_df = pd.DataFrame(meta_rows)
    metadata_df.to_excel(writer, sheet_name="Metadata", index=False)
    print("✓ Sheet: Metadata")

    writer.close()

    apply_excel_formatting(EXCEL_OUTPUT)
    print(f"\n✓ Excel workbook saved: {EXCEL_OUTPUT}")
    return EXCEL_OUTPUT


def apply_excel_formatting(filepath: Path):
    """Simple header formatting + auto column widths."""
    if not HAS_OPENPYXL or not filepath.exists():
        return

    wb = load_workbook(filepath)

    header_fill = PatternFill(
        start_color="1F4E78", end_color="1F4E78", fill_type="solid"
    )
    header_font = Font(color="FFFFFF", bold=True, size=11)

    for sheet in wb.worksheets:
        # Header row
        if sheet.max_row >= 1:
            for cell in sheet[1]:
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal="center", vertical="center")

        # Column widths
        for column in sheet.columns:
            max_len = 0
            col_letter = column[0].column_letter
            for cell in column:
                try:
                    val = str(cell.value) if cell.value is not None else ""
                    if len(val) > max_len:
                        max_len = len(val)
                except Exception:
                    continue
            sheet.column_dimensions[col_letter].width = min(max_len + 2, 50)

    wb.save(filepath)


# ======================================================================
# POWERPOINT DASHBOARD
# ======================================================================
def create_powerpoint_dashboard(data, charts):
    if not HAS_PPTX:
        print("\n✗ Skipping PowerPoint generation (python-pptx not installed).")
        return None

    print("\n" + "=" * 80)
    print("📊 GENERATING POWERPOINT DASHBOARD")
    print("=" * 80)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_slide(slide, "PoF3 Analiz Raporu", datetime.now().strftime("%d %B %Y"))
    print("✓ Slide 1: Title")

    # Slide 2: Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_executive_summary(slide, data)
    print("✓ Slide 2: Yönetici Özeti")

    # Slide 3: Equipment Distribution
    if "equipment_dist" in charts:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        add_chart_slide(slide, "Ekipman Tip Dağılımı", charts["equipment_dist"])
        print("✓ Slide 3: Equipment Distribution")

    # Slide 4: Chronic Distribution
    if "chronic_dist" in charts:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        add_chart_slide(slide, "Kronik Ekipman Analizi (IEEE 1366)", charts["chronic_dist"])
        print("✓ Slide 4: Chronic Distribution")

    # Slide 5: PoF by Horizon
    if "pof_horizons" in charts:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        add_chart_slide(slide, "PoF Tahminleri (Çok Ufuklu, Ensemble Öncelikli)", charts["pof_horizons"])
        print("✓ Slide 5: PoF by Horizon")

    # Slide 6: Feature Importance
    if "feature_imp" in charts:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        add_chart_slide(slide, "Özellik Önem Sıralaması (SHAP)", charts["feature_imp"])
        print("✓ Slide 6: Feature Importance")

    # Slide 7: Temporal Trends
    if "fault_trends" in charts:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        add_chart_slide(slide, "Arıza Trend Analizi", charts["fault_trends"])
        print("✓ Slide 7: Temporal Trends")

    prs.save(PPTX_OUTPUT)
    print(f"\n✓ PowerPoint dashboard saved: {PPTX_OUTPUT}")
    return PPTX_OUTPUT


def add_title_slide(slide, title, subtitle):
    box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    frame = box.text_frame
    p = frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.7))
    sub_frame = sub_box.text_frame
    p2 = sub_frame.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(22)
    p2.alignment = PP_ALIGN.CENTER


def add_executive_summary(slide, data):
    title = slide.shapes.title
    title.text = "Yönetici Özeti"

    base_df = (
        data["risk_master"]
        if not data["risk_master"].empty
        else (data["features"] if not data["features"].empty else data["equipment"])
    )
    total_eq = len(base_df) if base_df is not None and not base_df.empty else 0

    chronic_eq = 0
    if not data["chronic"].empty and "Kronik_Seviye_Max" in data["chronic"].columns:
        chronic_eq = (data["chronic"]["Kronik_Seviye_Max"] != "NORMAL").sum()

    meta = data.get("meta", {})
    data_start = meta.get("DATA_START_DATE") or "N/A"
    data_end = meta.get("DATA_END_DATE") or "N/A"

    # PoF model info is high level (no hard-coded metrics)
    text = f"""
• Toplam Ekipman: {total_eq:,}
• Kronik Ekipman (KRITIK/YUKSEK/ORTA/IZLEME): {chronic_eq:,} ({(100*chronic_eq/total_eq):.1f}%)
• Veri Aralığı: {data_start} – {data_end}
• Kullanılan Modeller:
    - Cox PH (kalibre, ekipman bazlı)
    - Weibull AFT (parametrik, yaşlanma analizi)
    - Random Survival Forest (RSF)
    - ML (XGBoost, kaçak etkisiz)
    - ENSEMBLE (Cox + Weibull + RSF + ML) → ÖNERİLEN PoF
"""

    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5)

    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True

    p = frame.paragraphs[0]
    p.text = text.strip()
    p.font.size = Pt(18)


def add_chart_slide(slide, title_text, chart_path: Path):
    title = slide.shapes.title
    title.text = title_text

    left = Inches(0.5)
    top = Inches(1.5)
    height = Inches(5.5)

    slide.shapes.add_picture(str(chart_path), left, top, height=height)


# ======================================================================
# MAIN
# ======================================================================
def main():
    print("\n" + "=" * 80)
    print("🚀 PoF3 DELIVERABLES GENERATOR v3 (Cox + Weibull + RSF + ML + ENSEMBLE)")
    print("=" * 80)
    print(f"Timestamp: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

    data = load_all_data()
    charts = create_visualizations(data)

    excel_file = create_excel_workbook(data, charts)
    pptx_file = create_powerpoint_dashboard(data, charts)

    print("\n" + "=" * 80)
    print("✅ DELIVERABLES GENERATION COMPLETED")
    print("=" * 80)

    if excel_file is not None and excel_file.exists():
        print(f"📊 Excel: {excel_file}")
        print(f"   Size: {excel_file.stat().st_size / 1024:.1f} KB")

    if pptx_file is not None and pptx_file.exists():
        print(f"📊 PowerPoint: {pptx_file}")
        print(f"   Size: {pptx_file.stat().st_size / 1024:.1f} KB")

    print(f"\n📁 Visualizations folder: {VIS_DIR}")
    print(f"   Charts generated: {len(charts)}")

    print("\n" + "=" * 80)
    print("🎉 All PoF3 v3 deliverables ready for stakeholders.")
    print("=" * 80)


if __name__ == "__main__":
    main()
